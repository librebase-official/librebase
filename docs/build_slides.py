#!/usr/bin/env python3
"""
LibreBase branded slide deck (python-pptx) — dark theme, teal accent, Orbitron/Space Grotesk/IBM Plex Mono.
Rebuilds the 12-slide market-research deck ON BRAND — same fonts and colors as librebase.xyz.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ═══════════════════════════════════════════
#  LIBREBASE BRAND TOKENS (from landing CSS + brand docs)
# ═══════════════════════════════════════════
# Dark theme (from globals.css :root + .lb-landing)
INK      = RGBColor(0x07, 0x10, 0x14)   # --lb-ink   bg
PANEL    = RGBColor(0x0E, 0x1A, 0x1C)   # --lb-panel
FOG      = RGBColor(0x9B, 0xB0, 0xAA)   # --lb-fog   muted text
PAPER    = RGBColor(0xE7, 0xF2, 0xEC)   # --lb-paper main text
SIGNAL   = RGBColor(0x2F, 0xD4, 0xC2)   # --lb-signal teal accent
SIGNAL_D = RGBColor(0x1F, 0xA8, 0x9A)   # --lb-signal-dim
WARN     = RGBColor(0xE8, 0xA8, 0x38)   # --lb-warn amber
LINE     = RGBColor(0xE7, 0xF2, 0xEC)   # --lb-line  (but at 0.12 alpha in CSS)
LINE_DIM = RGBColor(0x2D, 0x3A, 0x4D)   # border color from app shell
BORDER   = RGBColor(0x2D, 0x3A, 0x4D)   # --border
SURFACE  = RGBColor(0x1A, 0x23, 0x32)   # --surface

# Light theme tokens (from brand docs — used for PPTX text contrast)
L_BG     = RGBColor(0xFF, 0xFF, 0xFF)
L_TEXT   = RGBColor(0x0A, 0x0A, 0x0A)
L_ACCENT = RGBColor(0x49, 0x5D, 0x36)   # #495d36 green accent
L_MUTED  = RGBColor(0x52, 0x52, 0x5B)
L_BORDER = RGBColor(0xE4, 0xE4, 0xE7)
L_BG_MUT = RGBColor(0xF4, 0xF4, 0xF5)

# We use the dark cinematic theme for slides (matches librebase.xyz landing)
BG       = INK
TEXT     = PAPER
ACCENT   = SIGNAL
ACCENT_D = SIGNAL_D
MUTED    = FOG
PANEL_BG = PANEL
BORDER_C = BORDER
WARN_C   = WARN
RADIUS   = 8  # px (from --radius: 8px)

# Font files (downloaded from Google Fonts, same as website)
FONT_DIR = "/Users/julian/Downloads/librebase-fonts"
ORBITRON   = os.path.join(FONT_DIR, "Orbitron.ttf")
SPACE_GROT = os.path.join(FONT_DIR, "SpaceGrotesk.ttf")
IBM_PLEX   = os.path.join(FONT_DIR, "IBMPlexMono-Regular.ttf")

OUTPUT = "/Users/julian/Documents/coding-projects/librebase/docs/librebase-market-research-slides.pptx"

# ═══════════════════════════════════════════
#  HELPERS
# ═══════════════════════════════════════════

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def add_slide():
    return prs.slides.add_slide(BLANK)

def add_rect(slide, left, top, width, height, fill=None, line=None,
             line_w=None, shadow=False, rounded=False, border=None):
    # Support 'border' as alias for 'line'
    if border is not None and line is None:
        line = border
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line
        if line_w:
            shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    if rounded:
        # Set corner radius via XML
        try:
            sp_pr = shp._element.spPr
            # 8px on a shape is roughly 8/914400 inches — use adjustment
            adj = sp_pr.find(qn('a:prstGeom'))
            if adj is not None:
                # prst=roundedRect uses adj for corner radius (0..1, default 1/10)
                # We want ~8px which on a typical shape is ~0.089
                pass
        except Exception:
            pass
    return shp

def add_text(slide, text, left, top, width, height,
             font_size=14, color=TEXT, bold=False, font_name=SPACE_GROT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.4,
             italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    # Set East Asian / complex script font too
    for rPr in p._p.findall(qn('a:rPr')):
        rPr.set('altLang', 'en-US')
    p.alignment = align
    p.space_after = Pt(0)
    p.line_spacing = line_spacing
    # HACK: avoid broken relationship registration that breaks save()
    # (librebase-fonts are just for rendering; the .pptx stores font names,
    #  PowerPoint/Keynote handles font substitution if needed)
    return tf

def add_multi_text(slide, lines, left, top, width, height,
                   font_size=12, color=TEXT, font_name=SPACE_GROT,
                   line_spacing=1.5, align=PP_ALIGN.LEFT, bold_first=False):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(2)
        p.line_spacing = line_spacing
        if bold_first and i == 0:
            p.font.bold = True
    return tf

def add_bullet_list(slide, items, left, top, width, height,
                    font_size=12, color=TEXT, font_name=SPACE_GROT,
                    bullet_color=ACCENT, line_spacing=1.5, item_gap=6):
    """Bulleted list with accent-colored bullet markers."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    total_height = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(item_gap)
        p.line_spacing = line_spacing
        p.level = 0
        # Add bullet character manually via run prefix (easier than XML bullets)
        run_prefix = p.runs[0] if p.runs else None
        if not p.runs:
            run = p.add_run()
            run.text = "•  "
            run.font.size = Pt(font_size)
            run.font.color.rgb = bullet_color
            run.font.name = font_name
            run.font.bold = True
            p.text = item
        total_height += font_size * line_spacing + item_gap
    return tf

def add_card(slide, left, top, width, height, fill=PANEL_BG, border=BORDER_C, line_w=1):
    """A card panel — rounded rect with subtle border."""
    return add_rect(slide, left, top, width, height,
                    fill=fill, line=border, line_w=line_w, rounded=True)

def add_section_header(slide, title, subtitle=None, num=None):
    """Consistent section header used across most slides."""
    # Top accent bar
    add_rect(slide, 0, 0, 13.333, 0.06, fill=ACCENT)
    # Title
    add_text(slide, title, 0.55, 0.22, 12.2, 0.55,
             font_size=24, color=TEXT, bold=True, font_name=ORBITRON)
    # Thin rule under title
    add_rect(slide, 0.55, 0.72, 3.5, 0.03, fill=ACCENT_D)
    if subtitle:
        add_text(slide, subtitle, 0.55, 0.76, 11.5, 0.3,
                 font_size=10, color=MUTED, font_name=SPACE_GROT)
    if num:
        add_text(slide, str(num), 12.7, 7.12, 0.5, 0.25,
                 font_size=8, color=MUTED, align=PP_ALIGN.RIGHT)

def add_footer(slide, num=None):
    """Standard footer bar."""
    # Bottom line
    add_rect(slide, 0, 7.15, 13.333, 0.02, fill=BORDER_C)
    # Footer text
    add_text(slide, "LibreBase  •  librebase.xyz  |  Confidential — Internal Research",
             0.4, 7.18, 10.5, 0.25, font_size=7.5, color=MUTED)
    if num is not None:
        add_text(slide, str(num), 12.5, 7.18, 0.7, 0.25,
                 font_size=7.5, color=MUTED, align=PP_ALIGN.RIGHT)

def add_teal_accent_bar(slide, top=0, height=0.055):
    """Thin teal accent strip at top."""
    add_rect(slide, 0, top, 13.333, height, fill=ACCENT)

def add_logo_mark(slide, left, top, size=0.75):
    """LibreBase wordmark-style logo: 'Libre' + teal 'base'."""
    # "Libre" in Orbitron, "base" in teal
    add_text(slide, "Libre", left, top, size * 1.6, size * 0.8,
             font_size=18 * size / 0.75, color=TEXT, bold=True, font_name=ORBITRON)
    add_text(slide, "base", left + size * 0.95, top, size * 0.9, size * 0.8,
             font_size=18 * size / 0.75, color=ACCENT, bold=True, font_name=ORBITRON)

# ═══════════════════════════════════════════
#  SLIDE 1 — TITLE
# ═══════════════════════════════════════════

s = add_slide()
# Full dark bg
add_rect(s, 0, 0, 13.333, 7.5, fill=BG)

# Teal accent bar top
add_rect(s, 0, 0, 13.333, 0.06, fill=ACCENT)

# Background gradient orbs (matching landing page hero)
orb1 = add_rect(s, -1.5, -1.0, 4.5, 4.5, fill=ACCENT, shadow=False, rounded=True)
orb1.fill.fore_color.rgb = RGBColor(0x2F, 0xD4, 0xC2)
# Make it subtle with transparency via XML
try:
    solidFill = orb1.fill._fill
    srgb = solidFill.find(qn('a:solidFill')).find(qn('a:srgbClr'))
    if srgb is not None:
        alpha = srgb.makeelement(qn('a:alpha'), {'val': '12000'})  # 12% opacity
        srgb.append(alpha)
except Exception:
    pass

# Radial glow effect via a second blurred orb (simulated with large rounded rect)
orb2 = add_rect(s, 5.0, 2.5, 6.0, 6.0, fill=ACCENT_D, shadow=False, rounded=True)
try:
    solidFill2 = orb2.fill._fill
    srgb2 = solidFill2.find(qn('a:solidFill')).find(qn('a:srgbClr'))
    if srgb2 is not None:
        alpha2 = srgb2.makeelement(qn('a:alpha'), {'val': '6000'})  # 6% opacity
        srgb2.append(alpha2)
except Exception:
    pass

# Grid overlay (subtle lines matching landing page)
for i in range(0, 14, 2):
    add_rect(s, i * 0.95, 0, 0.005, 7.5, fill=RGBColor(0xE7, 0xF2, 0xEC),
             line=None, shadow=False)
for i in range(0, 8, 2):
    add_rect(s, 0, i * 0.95, 13.333, 0.005, fill=RGBColor(0xE7, 0xF2, 0xEC),
             line=None, shadow=False)
# Fade grid with transparency
for shape in s.shapes:
    try:
        sp_pr = shape._element.spPr
        solidFill = sp_pr.find(qn('a:solidFill'))
        if solidFill is not None:
            srgb = solidFill.find(qn('a:srgbClr'))
            if srgb is not None:
                clr_val = srgb.get('val')
                if clr_val and clr_val.startswith('E7F2EC'):
                    alpha = srgb.makeelement(qn('a:alpha'), {'val': '1500'})
                    srgb.append(alpha)
    except Exception:
        pass

# Logo wordmark
add_logo_mark(s, 0.7, 1.2, size=0.95)

# Title
add_text(s, "Market Intelligence",
         0.7, 2.15, 11.5, 0.9, font_size=38, color=TEXT, bold=True,
         font_name=ORBITRON, line_spacing=1.1)
add_text(s, "Sample Library & AI Music Tools Landscape",
         0.7, 3.0, 11.5, 0.45, font_size=18, color=ACCENT, font_name=SPACE_GROT,
         line_spacing=1.3)

# Teal rule
add_rect(s, 0.7, 3.55, 2.5, 0.04, fill=ACCENT)

# Meta info
add_multi_text(s, [
    "INTERNAL RESEARCH BRIEF  •  FOR SALES USE ONLY",
    "Prepared for: Sales Representative — Internal Use  |  Date: August 16–18, 2026",
    "Classification: Confidential"
], 0.7, 3.8, 11.5, 1.2, font_size=10, color=MUTED, font_name=SPACE_GROT,
   line_spacing=1.6)

# Bottom: brand URL
add_text(s, "librebase.xyz", 0.7, 6.6, 4, 0.3,
         font_size=11, color=ACCENT, bold=True, font_name=ORBITRON)

# Footer
add_footer(s, 1)

# ═══════════════════════════════════════════
# SLIDE 2 — EXECUTIVE SUMMARY
# ═══════════════════════════════════════════

s = add_slide()
add_rect(s, 0, 0, 13.333, 7.5, fill=BG)
add_teal_accent_bar(s)
add_section_header(s, "Executive Summary", num=2)
add_footer(s, 2)

# Left: Key insight panel
add_card(s, 0.55, 0.95, 5.8, 3.0, fill=PANEL_BG, border=BORDER_C, line_w=1)
# Panel header strip
add_rect(s, 0.55, 0.95, 5.8, 0.35, fill=ACCENT_D)
add_text(s, "THE ONE-LINE PITCH", 0.7, 0.98, 5.4, 0.3,
         font_size=10, color=INK, bold=True, font_name=ORBITRON)
add_multi_text(s, [
    "The sample market is $1.8B and growing at 8.9–9.5% CAGR.",
    "",
    "Splice (22% share, 6M+ users) dominates volume;",
    "Loopcloud owns DAW integration; Tracklib monopolises",
    "legal clearance. But every incumbent charges subscription",
    "rents that producers hate, hides quality behind credit",
    "systems, and lacks any all-in-one AI-native workflow.",
    "",
    "The gap is real — and wide."
], 0.7, 1.38, 5.4, 2.4, font_size=11, color=TEXT, font_name=SPACE_GROT,
   line_spacing=1.35)

# Right: Why it matters
add_text(s, "Why This Matters for Sales", 6.65, 0.95, 6.2, 0.4,
         font_size=18, color=TEXT, bold=True, font_name=ORBITRON)

bullets_right = [
    ("Mid-market is underserved",
     "Enterprise tools (Sphera, Enablon) cost six figures; consumer subs ($8–40/mo) feel "
     "predatory. The $150–300 'professional-lite' tier is empty."),
    ("AI is the wedge, not the product",
     "Every competitor's AI is bolted-on and mediocre. The gap: unify generation + separation "
     "+ editing + MIDI in one plugin."),
    ("Perpetual ownership = emotional hook",
     "'Cancel = lose everything' is #1 pain point across Splice, Loopcloud, Output, EastWest. "
     "A hybrid ownership model is unclaimed."),
    ("Consolidation breeds vulnerability",
     "Splice owns Spitfire, Output, Sample Magic. Users are anxious; a vendor-neutral, open "
     "platform is a natural trust alternative."),
]

y = 1.55
for title, desc in bullets_right:
    # Accent bullet dot
    add_rect(s, 6.65, y + 0.06, 0.1, 0.1, fill=ACCENT)
    add_text(s, title, 6.85, y - 0.04, 5.8, 0.28,
             font_size=11, color=ACCENT, bold=True, font_name=ORBITRON)
    add_text(s, desc, 6.85, y + 0.22, 5.8, 0.7, font_size=9.5, color=TEXT,
             line_spacing=1.25)
    y += 1.25


# ═══════════════════════════════════════════
# SLIDE 3 — MARKET SNAPSHOT
# ═══════════════════════════════════════════

s = add_slide()
add_rect(s, 0, 0, 13.333, 7.5, fill=BG)
add_teal_accent_bar(s)
add_section_header(s, "Market Snapshot", num=3)
add_footer(s, 3)

rows = [
    ["Market size (2025)", "$1.8B", "Large enough to matter; not saturated"],
    ["Projected (2034)", "$3.9–4.6B", "~9% CAGR — sustained growth, not a bubble"],
    ["Splice market share", "~22%", "Leading, but not a monopoly"],
    ["Splice registered users", "6M+", "Huge installed base — churn opportunity"],
    ["Splice 2024 downloads", "~350M", "Scale, but 'generic' quality complaints"],
    ["Spotify-for-samples sentiment", "Broken", "Ownership is the untapped wedge"],
    ["AI sample gen market", "Fragmented", "No dominant all-in-one — white space"],
    ["SDS/Compliance adj. market", "$1.34B → $2.49B (14% CAGR)",
     "Parallel vertical, same open-data angle"],
]

# Table header
add_text(s, "METRIC", 0.55, 1.0, 3.3, 0.3, font_size=9, color=ACCENT,
         bold=True, font_name=ORBITRON)
add_text(s, "VALUE", 3.9, 1.0, 2.3, 0.3, font_size=9, color=ACCENT,
         bold=True, font_name=ORBITRON)
add_text(s, "SALES SIGNAL", 6.3, 1.0, 6.0, 0.3, font_size=9, color=ACCENT,
         bold=True, font_name=ORBITRON)
add_rect(s, 0.55, 1.25, 12.2, 0.02, fill=ACCENT_D)

y = 1.35
for i, (metric, value, signal) in enumerate(rows):
    bg_c = PANEL_BG if i % 2 == 0 else RGBColor(0x0A, 0x15, 0x18)
    add_rect(s, 0.55, y, 12.2, 0.52, fill=bg_c, line=BORDER_C, line_w=0.5, rounded=True)
    add_text(s, metric, 0.68, y + 0.06, 3.1, 0.38, font_size=9.5, color=TEXT)
    add_text(s, value, 4.0, y + 0.06, 2.1, 0.38, font_size=9.5, color=ACCENT,
             bold=True)
    add_text(s, signal, 6.4, y + 0.06, 6.1, 0.38, font_size=9.5, color=MUTED,
             line_spacing=1.2)
    y += 0.55


# ═══════════════════════════════════════════
# SLIDE 4 — COMPETITIVE: 3 COLUMNS
# ═══════════════════════════════════════════

s = add_slide()
add_rect(s, 0, 0, 13.333, 7.5, fill=BG)
add_teal_accent_bar(s)
add_section_header(s, "Competitive Landscape — Subscription Platforms", num=4)
add_footer(s, 4)

col_data = [
    {"title": "Splice", "sub": "The Incumbent Leader", "color": ACCENT,
     "items": [
        ("Library", "3–4M+ samples"),
        ("Pricing", "Sounds+ $12.99/mo | Creator $19.99/mo\nCreator+ $39.99/mo | Annual promo: $120/yr"),
        ("Strengths", "Largest library, brand recognition, AI search,\nRent-to-Own, Splice Studio. Acquired Spitfire\n($50M) + Output + Sample Magic — consolidation king."),
        ("Weaknesses", "Credit system restrictive; price hikes\n(Sounds 1000 → $34.99/mo); lose access\non cancel; credits expire. 80% 'generic'."),
     ]},
    {"title": "Loopcloud", "sub": "DAW Integration Leader", "color": ACCENT_D,
     "items": [
        ("Library", "4M+ samples"),
        ("Pricing", "Artist $7.99/mo | Studio $11.99/mo\n(most popular) | Professional tier"),
        ("Strengths", "Deepest DAW integration (VST/AAX/AU),\nno credit limits (points-based packs),\nbest for house/techno, cloud storage\n(10–50GB+), AI browsing, time-stretch in-plugin."),
        ("Weaknesses", "Internet required for cloud features;\nplugin learning curve;\nlibrary overlap with Splice."),
     ]},
    {"title": "Tracklib", "sub": "Legal Clearance Monopoly", "color": WARN_C,
     "items": [
        ("Library", "Original songs (not royalty-free loops)"),
        ("Pricing", "Lite $8.99/mo | Premium $13.99/mo\nMax $19.99/mo"),
        ("Strengths", "One-of-a-kind — real records legally\nsampleable; used by Kendrick Lamar\nproducers; no upfront fees;\ndrag-to-DAW desktop app."),
        ("Weaknesses", "No free account; Trustpilot 2.0/5;\nconfusing tier structure (A/B/C\nclearance); limited royalty-free genres."),
     ]},
]

col_width = 3.9
x_positions = [0.45, 4.65, 8.85]

for idx, col in enumerate(col_data):
    x = x_positions[idx]
    # Header bar
    add_rect(s, x, 0.95, col_width, 0.62, fill=col["color"], rounded=True)
    add_text(s, col["title"], x + 0.15, 0.97, col_width - 0.3, 0.3,
             font_size=14, color=INK if col["color"] != WARN_C else RGBColor(0x1a,0x1a,0x1a),
             bold=True, font_name=ORBITRON)
    add_text(s, col["sub"], x + 0.15, 1.23, col_width - 0.3, 0.25,
             font_size=8.5, color=INK if col["color"] != WARN_C else RGBColor(0x3a,0x3a,0x3a),
             font_name=SPACE_GROT)
    y = 1.72
    for label, content in col["items"]:
        add_text(s, label, x + 0.15, y, col_width - 0.3, 0.2,
                 font_size=8.5, color=col["color"], bold=True, font_name=ORBITRON)
        y += 0.19
        for line in content.split("\n"):
            add_text(s, line, x + 0.25, y, col_width - 0.4, 0.18,
                     font_size=8, color=TEXT, line_spacing=1.15)
            y += 0.17
        y += 0.15


# ═══════════════════════════════════════════
# SLIDE 5 — THE REST + AI MATRIX
# ═══════════════════════════════════════════

s = add_slide()
add_rect(s, 0, 0, 13.333, 7.5, fill=BG)
add_teal_accent_bar(s)
add_section_header(s, "The Rest of the Market + AI Quality Matrix", num=5)
add_footer(s, 5)

# Left: Other Players table
add_text(s, "OTHER PLAYERS", 0.55, 0.82, 5.5, 0.3,
         font_size=11, color=ACCENT, bold=True, font_name=ORBITRON)
add_rect(s, 0.55, 1.08, 1.2, 0.02, fill=ACCENT_D)

other_rows = [
    ["Noiiz", "Unlimited", "$5–$10/mo", "Small lib; Tp 2.2/5"],
    ["BandLab", "Social", "Free/$8–$16", "Dark horse — Gen Z pull"],
    ["Output Arcade", "Sampler", "$9.99–$12.99", "Sub-only; Splice-owned"],
    ["Output One", "Full suite", "$14.99/mo", "Locked ecosystem"],
    ["LANDR Studio", "Bundled", "$8.25–$19.99", "Samples + mastering"],
    ["Sample Focus", "Community", "Free/Premium", "Community-curated"],
    ["NI Sounds.com", "Sound design", "$9.99/mo", "Ambient/experimental"],
    ["Waves ILLUGEN", "Text-to-sound", "One-time/sub", "MusicTech rated 5/10"],
]

y = 1.18
add_text(s, "Platform", 0.55, y, 1.5, 0.22, font_size=8, color=ACCENT,
         bold=True, font_name=ORBITRON)
add_text(s, "Model", 2.1, y, 1.1, 0.22, font_size=8, color=ACCENT,
         bold=True, font_name=ORBITRON)
add_text(s, "Price", 3.25, y, 1.2, 0.22, font_size=8, color=ACCENT,
         bold=True, font_name=ORBITRON)
add_text(s, "Watch", 4.5, y, 1.6, 0.22, font_size=8, color=ACCENT,
         bold=True, font_name=ORBITRON)
add_rect(s, 0.55, y + 0.2, 5.55, 0.01, fill=ACCENT_D)
y += 0.24

for i, row in enumerate(other_rows):
    bg_c = PANEL_BG if i % 2 == 0 else RGBColor(0x0A, 0x15, 0x18)
    add_rect(s, 0.55, y, 5.55, 0.27, fill=bg_c, line=BORDER_C, line_w=0.3, rounded=True)
    add_text(s, row[0], 0.62, y + 0.02, 1.45, 0.22, font_size=7.5, color=TEXT)
    add_text(s, row[1], 2.12, y + 0.02, 1.05, 0.22, font_size=7.5, color=MUTED)
    add_text(s, row[2], 3.25, y + 0.02, 1.2, 0.22, font_size=7.5, color=TEXT)
    add_text(s, row[3], 4.5, y + 0.02, 1.5, 0.22, font_size=7.5, color=MUTED)
    y += 0.29

# Right: AI Quality Matrix
add_text(s, "AI TOOLS — FEATURE MATRIX", 6.45, 0.82, 6.4, 0.3,
         font_size=11, color=ACCENT, bold=True, font_name=ORBITRON)
add_rect(s, 6.45, 1.08, 1.8, 0.02, fill=ACCENT_D)

ai_rows = [
    ("Samplab (RIP)", "✓", "✓", "✓", "✓", "VST3", "$12/mo"),
    ("RipX DAW", "—", "✓", "✓", "✓", "Std", "$60–$149"),
    ("Waves ILLUGEN", "✓", "—", "—", "—", "Std", "$8–$20/mo"),
    ("Text-to-Sample", "✓", "—", "—", "—", "VST3/AU", "Pay/use"),
    ("Magenta (Google)", "✓", "—", "—", "—", "Open", "Free"),
    ("Splice Create", "✓", "—", "—", "—", "Web", "$13/mo"),
    ("Loudly", "✓", "✓", "—", "—", "Web", "Freemium"),
    ("NeuralNote", "—", "—", "✓", "—", "VST3/AU", "Free"),
    ("Emergence Audio", "—", "—", "—", "—", "Kontakt", "$29–$1,299"),
    ("Adobe Podcast", "—", "✓", "—", "—", "Web", "Free/Prem"),
    ("BandLab", "✓", "—", "—", "—", "Web", "Free/$8–16"),
]

col_ws = [1.45, 0.4, 0.4, 0.55, 0.5, 0.65, 0.75]
x_start = 6.45
y = 1.18

headers = ["Tool", "Gen", "Sep", "A2MIDI", "Edit", "Plugin", "Price"]
for hi, h in enumerate(headers):
    add_text(s, h, x_start, y, col_ws[hi], 0.22, font_size=7.5, color=ACCENT,
             bold=True, font_name=ORBITRON, align=PP_ALIGN.CENTER)
    x_start += col_ws[hi]
x_start = 6.45
add_rect(s, 6.45, y + 0.2, sum(col_ws), 0.01, fill=ACCENT_D)
y += 0.24

for i, row in enumerate(ai_rows):
    bg_c = PANEL_BG if i % 2 == 0 else RGBColor(0x0A, 0x15, 0x18)
    add_rect(s, 6.45, y, sum(col_ws), 0.24, fill=bg_c, line=BORDER_C, line_w=0.3, rounded=True)
    x_c = 6.45
    for ci, cell in enumerate(row):
        if cell == "✓":
            c = ACCENT
            b = True
        elif cell == "—":
            c = RGBColor(0x6B, 0x72, 0x80)
            b = False
        elif ci == 0:
            c = TEXT
            b = True
        elif ci == len(row) - 1:
            c = MUTED
            b = False
        else:
            c = TEXT
            b = False
        add_text(s, cell, x_c + 0.02, y + 0.01, col_ws[ci] - 0.04, 0.2,
                 font_size=7, color=c, bold=b, align=PP_ALIGN.CENTER)
        x_c += col_ws[ci]
    y += 0.26


# ═══════════════════════════════════════════
# SLIDE 6 — AI PLAYERS + GAPS
# ═══════════════════════════════════════════

s = add_slide()
add_rect(s, 0, 0, 13.333, 7.5, fill=BG)
add_teal_accent_bar(s)
add_section_header(s, "AI Tools — Primary Players & Gaps", num=6)
add_footer(s, 6)

players = [
    ("Samplab — SHUTTING DOWN Sept 17, 2026",
     "Polyphonic audio-to-MIDI, stem separation, chord detection, note-level editing, TextToSample. "
     "Only tool doing note-level editing inside polyphonic audio seamlessly. No successor. "
     "Gap left behind: biggest hole in the market.",
     ACCENT),
    ("Waves ILLUGEN — Text-to-Sound",
     "Text → one-shots, loops, SFX. Standalone desktop app only (no plugin). Credit-based. "
     "New; MusicTech rated 5/10. Quality concerns.",
     ACCENT_D),
    ("Text-to-Sample.com — New Entrant",
     "Text → audio (5–30s). Web + VST3/AU plugin. Credit-based, pay-as-you-go. No vocals. "
     "Short max length. Unproven.",
     ACCENT_D),
    ("Splice Create — April 2026 Launch",
     "Three generative AI tools for reshaping Splice library sounds. Text-to-sample within "
     "Splice ecosystem. Credit-locked. Quality unproven.",
     ACCENT),
    ("Google Magenta — Open Source, Free",
     "Magenta RealTime 2 (800M param transformer, 190k hrs stock music, June 2025). "
     "Magenta Studio for Ableton. DDSP-VST, NSynth. Requires technical setup. Research-grade.",
     ACCENT_D),
    ("RipX DAW — Samplab Replacement (Partial)",
     "6+ stem AI separation, note-level extraction/editing, instrument replacement. "
     "One-time purchase ($60–$149). Steep learning curve. Not a generator — separation + editing only.",
     ACCENT_D),
]

y = 0.98
for title, desc, color in players:
    add_rect(s, 0.55, y + 0.04, 0.08, 0.7, fill=color)
    add_text(s, title, 0.78, y - 0.03, 12.2, 0.25, font_size=11, color=color,
             bold=True, font_name=ORBITRON)
    add_text(s, desc, 0.78, y + 0.2, 12.2, 0.42, font_size=9, color=TEXT,
             line_spacing=1.2)
    y += 0.75

# Gaps callout
add_rect(s, 0.55, y + 0.08, 12.2, 0.55, fill=PANEL_BG, border=ACCENT, line_w=1, rounded=True)
add_text(s, "KEY GAPS", 0.7, y + 0.1, 2.0, 0.22, font_size=10, color=ACCENT,
         bold=True, font_name=ORBITRON)
add_text(s, "Polyphonic note editing void  |  No all-in-one tool  |  AI vocals underdeveloped  |  "
             "Real-time tools absent  |  Plugin integration gap  |  Orchestral/world instruments poor",
         0.7, y + 0.32, 11.8, 0.22, font_size=9.5, color=TEXT, line_spacing=1.2)


# ═══════════════════════════════════════════
# SLIDE 7 — TRADITIONAL LIBRARIES MAP
# ═══════════════════════════════════════════

s = add_slide()
add_rect(s, 0, 0, 13.333, 7.5, fill=BG)
add_teal_accent_bar(s)
add_section_header(s, "Traditional Sample Library Competitors",
                   subtitle="Market: Orchestral, cinematic, traditional libraries  |  ~$1.8B (2025) → $3.9B by 2034",
                   num=7)
add_footer(s, 7)

trad_rows = [
    ["Spitfire Audio", "BBC SO Pro €999", "LABS+ £12.99/mo", "LABS free", "SINE", "Expensive; Splice fears; bugs"],
    ["Native Instruments", "Komplete Ult. $1,799", "—", "—", "iLok", "Dongle hatred; insolvency rumors; heavy CPU"],
    ["Output", "—", "Arcade $12.99/mo", "—", "—", "Sub-only; trust (Splice-owned); mixed quality"],
    ["Heavyocity", "NOVO/Damage $399", "—", "—", "Kontakt", "Niche; expensive; Kontakt dep."],
    ["Cinesamples", "Musio $299–$399", "Musio $9.99/mo", "—", "Musio", "'Too good to be true'; player quality"],
    ["8Dio", "Adagio $499", "—", "—", "Kontakt", "QC issues; Trustpilot 3.1/8"],
    ["Embertone", "Joshua Bell $199", "—", "—", "Kontakt", "Niche (solo); small catalog"],
    ["Orch. Tools", "Berlin Max €1,399", "—", "Berlin Free", "SINE", "Very expensive; bugs; no Linux"],
    ["EastWest", "—", "ComposerCloud $19.99", "—", "Install ctr", "Worst portal software; sub fatigue"],
]

cws = [1.3, 1.4, 1.25, 1.0, 0.75, 2.3]
x0 = 0.55
y = 0.98
headers = ["Competitor", "Premium Anchor", "Subscription", "Free Entry", "DRM", "Core Pain"]
for hi, h in enumerate(headers):
    add_text(s, h, x0, y, cws[hi], 0.22, font_size=8, color=ACCENT,
             bold=True, font_name=ORBITRON)
    x0 += cws[hi]
x0 = 0.55
add_rect(s, 0.55, y + 0.2, sum(cws), 0.01, fill=ACCENT_D)
y += 0.24

for i, row in enumerate(trad_rows):
    bg_c = PANEL_BG if i % 2 == 0 else RGBColor(0x0A, 0x15, 0x18)
    add_rect(s, 0.55, y, sum(cws), 0.3, fill=bg_c, line=BORDER_C, line_w=0.3, rounded=True)
    x_c = 0.55
    for ci, cell in enumerate(row):
        col = ACCENT if ci == 0 else TEXT
        add_text(s, cell, x_c + 0.03, y + 0.02, cws[ci] - 0.06, 0.24,
                 font_size=7.5, color=col, bold=(ci == 0), line_spacing=1.15)
        x_c += cws[ci]
    y += 0.32


# ═══════════════════════════════════════════
# SLIDE 8 — PAIN POINTS
# ═══════════════════════════════════════════

s = add_slide()
add_rect(s, 0, 0, 13.333, 7.5, fill=BG)
add_teal_accent_bar(s)
add_section_header(s, "Industry Pain Points — The Sales Playbook", num=8)
add_footer(s, 8)

pains = [
    ("1. The 'Rental' Problem",
     "Cancel subscription → lose access to downloaded samples. Credits expire on cancellation. "
     "No perpetual ownership in most subscriptions.",
     "Sales angle: 'You keep what you download — forever.' Ownership as differentiator."),
    ("2. Credit/Point System Frustration",
     "Artificial currency obscures real cost; feels predatory. Hard to know cost-per-sample. "
     "Unused credits create anxiety (use-it-or-lose-it).",
     "Sales angle: Transparent pricing. No credits. Know what you pay."),
    ("3. Discovery & Quality Inconsistency",
     "80% of Splice samples criticised as generic. Poor metadata/tempo-key tagging on one-shots. "
     "Genre gaps: world music, non-Western instruments, niche genres.",
     "Sales angle: Curated quality over volume. Better metadata. Genre depth."),
    ("4. Technical Pain Points",
     "Missing samples when collaborating (Splice). DAW integration clunky (Splice) vs seamless "
     "(Loopcloud). Always-on internet required for cloud. Desktop apps resource-heavy.",
     "Sales angle: Offline-first. DAW-native plugin. No cloud dependency."),
    ("5. Pricing & Value Fatigue",
     "Multiple subscriptions add up ($10–40/mo each). Price hikes (Splice Sounds 1000 → $34.99/mo). "
     "Hard to justify for hobbyists. Annual commitments feel risky.",
     "Sales angle: One subscription. Flat pricing. No per-user tax."),
    ("6. Creator/Artist Compensation",
     "Producers report low per-download payouts on Splice. No transparency in royalty structures.",
     "Sales angle: Transparent per-download payouts. Creator profiles. Tip jars."),
]

y = 0.95
for title, desc, angle in pains:
    add_rect(s, 0.55, y + 0.04, 0.06, 0.72, fill=WARN_C)
    add_text(s, title, 0.75, y - 0.03, 3.5, 0.22, font_size=10.5, color=TEXT,
             bold=True, font_name=ORBITRON)
    add_text(s, desc, 0.75, y + 0.18, 6.5, 0.5, font_size=8.5, color=TEXT,
             line_spacing=1.2)
    add_text(s, angle, 7.5, y + 0.18, 5.3, 0.5, font_size=8.5, color=ACCENT,
             line_spacing=1.2)
    y += 0.82


# ═══════════════════════════════════════════
# SLIDE 9 — OPPORTUNITIES
# ═══════════════════════════════════════════

s = add_slide()
add_rect(s, 0, 0, 13.333, 7.5, fill=BG)
add_teal_accent_bar(s)
add_section_header(s, "Market Gaps & Opportunities", num=9)
add_footer(s, 9)

opps = [
    ("Compressed/AI-Assisted Libraries",
     "Pro sound without 100GB+ per library. AI-assisted sample generation could reduce size "
     "10–50×. No player leadership in this space."),
    ("Perpetual-Friendly Subscription Hybrid",
     "Spitfire's Splice rent-to-own is a start. Users want 'own after X payments' model. "
     "No major player offers true rent-to-own for all libraries."),
    ("Playerless/Web-Based Solution",
     "iLok + proprietary players = friction. Web-based auth or simple license key. "
     "Cross-platform (Linux underserved — zero major library supports it)."),
    ("Transparent Mid-Tier Pricing",
     "Gap between $30–50 singles and $500–1000 orchestras. $150–300 'professional lite' "
     "tier underserved. Cinesamples Musio ($199–399) closest but quality questioned."),
    ("All-in-One Subscription with Ownership",
     "ComposerCloud $19.99/mo is good value but no ownership. Users want 'subscribe 24 months, "
     "own forever' option. Output One $14.99/mo has same criticism."),
    ("Stability & Performance Focus",
     "SINE player bugs, Kontakt crashes are common complaints. 'Just works' positioning "
     "could win market share. Compressed libraries = less RAM/CPU strain."),
    ("Linux Support",
     "Zero major sample library supports Linux. Clear niche for someone to own."),
]

y = 0.95
for title, desc in opps:
    add_rect(s, 0.55, y + 0.04, 0.06, 0.6, fill=ACCENT)
    add_text(s, title, 0.75, y - 0.03, 12.0, 0.2, font_size=10.5, color=TEXT,
             bold=True, font_name=ORBITRON)
    add_text(s, desc, 0.75, y + 0.16, 12.0, 0.35, font_size=8.5, color=TEXT,
             line_spacing=1.2)
    y += 0.72


# ═══════════════════════════════════════════
# SLIDE 10 — STRATEGIC POSITIONING
# ═══════════════════════════════════════════

s = add_slide()
add_rect(s, 0, 0, 13.333, 7.5, fill=BG)
add_teal_accent_bar(s)
add_section_header(s, "Strategic Positioning — Where LibreBase Fits", num=10)
add_footer(s, 10)

# Open knowledge thesis card
add_card(s, 0.55, 0.92, 12.2, 1.0, fill=PANEL_BG, border=ACCENT, line_w=1)
add_text(s, "THE OPEN-KNOWLEDGE THESIS", 0.7, 0.94, 5.0, 0.25,
         font_size=10, color=ACCENT, bold=True, font_name=ORBITRON)
add_text(s, "LibreBase's core thesis — knowledge should be findable, attributable, and freely "
             "reusable — maps directly onto every pain point above: No lock-in, full attribution, "
             "open metadata, perpetual access. Download and own. Cancel and keep.",
         0.7, 1.18, 11.8, 0.65, font_size=10, color=TEXT, line_spacing=1.3)

# Positioning table
add_text(s, "CONCRETE SALES POSITIONS", 0.55, 2.1, 6.0, 0.28,
         font_size=12, color=ACCENT, bold=True, font_name=ORBITRON)
add_rect(s, 0.55, 2.34, 1.5, 0.02, fill=ACCENT_D)

pos_rows = [
    ("Open platform vs. Splice walled garden",
     "Splice owns Spitfire, Output, Sample Magic. LibreBase is vendor-neutral: your data, "
     "your formats, your choice."),
    ("Perpetual ownership vs. rental",
     "Cancel Splice and you lose 3–4M samples. With LibreBase, what you acquire is yours — "
     "no rental, no credit expiry."),
    ("Transparent pricing vs. credit systems",
     "Splice's credit model obscures real cost. LibreBase uses flat, predictable pricing — "
     "no artificial currency, no use-it-or-lose-it."),
    ("AI-native vs. bolt-on AI",
     "Splice Create is AI bolted onto a library. LibreBase's open approach means generation + "
     "separation + editing + MIDI in one workflow."),
    ("No DRM vs. iLok / SINE / Kontakt",
     "iLok is universally hated. SINE crashes. Kontakt eats RAM. LibreBase: no dongles, "
     "no proprietary players, no Linux exclusion."),
    ("Mid-market focus vs. enterprise blanks",
     "Sphera / Enablon cost six figures. Consumer subs feel predatory. LibreBase targets "
     "the $150–300 professional-lite gap nobody fills."),
    ("Linux support as wedge",
     "Zero major sample library supports Linux. LibreBase does — own a niche nobody else touches."),
    ("Creator economy vs. black-box royalties",
     "Splice producers report low payouts with no transparency. LibreBase makes attribution "
     "and compensation visible by default."),
]

y = 2.45
for i, (pos, hook) in enumerate(pos_rows):
    bg_c = PANEL_BG if i % 2 == 0 else RGBColor(0x0A, 0x15, 0x18)
    add_rect(s, 0.55, y, 12.2, 0.37, fill=bg_c, line=BORDER_C, line_w=0.3, rounded=True)
    add_text(s, pos, 0.65, y + 0.02, 2.8, 0.3, font_size=8, color=ACCENT,
             bold=True, line_spacing=1.15)
    add_text(s, hook, 3.6, y + 0.02, 9.0, 0.3, font_size=8, color=TEXT,
             line_spacing=1.15)
    y += 0.39


# ═══════════════════════════════════════════
# SLIDE 11 — KEY NUMBERS + OBJECTIONS
# ═══════════════════════════════════════════

s = add_slide()
add_rect(s, 0, 0, 13.333, 7.5, fill=BG)
add_teal_accent_bar(s)
add_section_header(s, "Key Numbers & Top Objections", num=11)
add_footer(s, 11)

# Left: Key Numbers
add_text(s, "KEY NUMBERS TO QUOTE", 0.55, 0.82, 5.5, 0.3,
         font_size=11, color=ACCENT, bold=True, font_name=ORBITRON)
add_rect(s, 0.55, 1.08, 1.3, 0.02, fill=ACCENT_D)

nums = [
    ("$1.8B", "Current sample library market (2025)"),
    ("8.9–9.5% CAGR", "Sustained growth through 2034"),
    ("22%", "Splice market share (leader, not monopoly)"),
    ("6M+", "Splice registered users (churn opportunity)"),
    ("$34.99/mo", "Splice Sounds 1000 (price hike flashpoint)"),
    ("Trustpilot scores", "Splice 4.3 | Output 3.7 | EastWest 3.4 | "
     "8Dio 3.1 | Heavyocity 2.5 | Tracklib 2.0 | Noiiz 2.2"),
    ("Sept 17, 2026", "Samplab shutdown — leaves polyphonic note-editing gap"),
]

y = 1.18
for num, ctx in nums:
    add_rect(s, 0.55, y + 0.04, 0.08, 0.32, fill=ACCENT)
    add_text(s, num, 0.75, y - 0.02, 1.8, 0.22, font_size=12, color=ACCENT,
             bold=True, font_name=ORBITRON)
    add_text(s, ctx, 2.6, y - 0.02, 3.6, 0.35, font_size=9, color=TEXT,
             line_spacing=1.2)
    y += 0.42

# Right: Objections
add_text(s, "TOP 5 OBJECTIONS & RESPONSES", 6.55, 0.82, 6.2, 0.3,
         font_size=11, color=ACCENT, bold=True, font_name=ORBITRON)
add_rect(s, 6.55, 1.08, 2.2, 0.02, fill=ACCENT_D)

objs = [
    ("'Splice has the biggest library'",
     "Yes — 3–4M+ samples, 22% share. But 80% call them generic, credits expire, "
     "cancel means losing everything. We offer ownership + curation."),
    ("'Loopcloud integrates with my DAW'",
     "Loopcloud's plugin is strong, but cloud-dependent and points-based. We're "
     "DAW-native, offline-capable, flat-priced — no credit maths."),
    ("'I've paid for Kontakt libraries'",
     "Kontakt's iLok DRM is the #1 industry complaint. We're playerless — no dongles, "
     "no proprietary lock-in. Your libraries stay yours."),
    ("'AI sample gen is gimmicky'",
     "ILLUGEN scored 5/10. Splice Create is unproven. But Samplab's shutdown left a real "
     "gap — the market wants generation + editing in one tool."),
    ("'Enterprise tools are more trustworthy'",
     "Sphera / Enablon cost $16K–$100K+ and are overkill under $500M revenue. Mid-market "
     "is underserved. We're built for AI-powered review without enterprise tax."),
]

y = 1.18
for i, (obj, resp) in enumerate(objs):
    bg_c = PANEL_BG if i % 2 == 0 else RGBColor(0x0A, 0x15, 0x18)
    add_rect(s, 6.55, y, 6.25, 0.65, fill=bg_c, line=BORDER_C, line_w=0.3, rounded=True)
    add_text(s, obj, 6.68, y + 0.03, 6.0, 0.2, font_size=8.5, color=ACCENT,
             bold=True, font_name=ORBITRON)
    add_text(s, resp, 6.68, y + 0.22, 6.0, 0.35, font_size=7.5, color=TEXT,
             line_spacing=1.2)
    y += 0.7


# ═══════════════════════════════════════════
# SLIDE 12 — SIGNALS + CLOSING
# ═══════════════════════════════════════════

s = add_slide()
add_rect(s, 0, 0, 13.333, 7.5, fill=BG)
add_teal_accent_bar(s)
add_section_header(s, "Signals to Monitor & Closing", num=12)
add_footer(s, 12)

signals = [
    "Splice price trajectory post-Spitfire / Output acquisitions — user churn sentiment",
    "BandLab's free-tier growth — Gen Z creator pull",
    "Any Google productisation of Magenta RealTime 2 — free AI undercutting paid tools",
    "Text-to-sample quality improvements — when does it become production-ready?",
    "Linux adoption in music production — growing but still zero competitor support",
]

add_text(s, "SIGNALS TO MONITOR", 0.55, 0.82, 5.0, 0.3,
         font_size=11, color=ACCENT, bold=True, font_name=ORBITRON)
add_rect(s, 0.55, 1.08, 1.3, 0.02, fill=ACCENT_D)

y = 1.18
for sig in signals:
    add_rect(s, 0.55, y + 0.04, 0.1, 0.1, fill=ACCENT)
    add_text(s, sig, 0.78, y, 12.0, 0.28, font_size=10, color=TEXT)
    y += 0.3

# Closing thesis card
y += 0.12
add_card(s, 0.55, y, 12.2, 2.0, fill=PANEL_BG, border=ACCENT, line_w=1)
add_text(s, "THE LIBREBASE PITCH", 0.7, y + 0.1, 4.0, 0.25,
         font_size=11, color=ACCENT, bold=True, font_name=ORBITRON)
add_multi_text(s, [
    "LibreBase is the open knowledge base and competitive intelligence platform.",
    "",
    "It's the tool you point at a messy pile of market research, competitor intel,",
    "structured data, and org knowledge and turn it into something a team can actually",
    "use — without handing the keys to a vendor.",
    "",
    "For the sample library market: an open, vendor-neutral, DRM-free, perpetually-owned",
    "alternative to Splice's walled garden. For the AI music tools market: a unifying",
    "platform where generation, separation, editing, and MIDI converge — not a crediting game.",
    "",
    "librebase.xyz  |  app.librebase.xyz  |  The open knowledge base, owned by the people who use it.",
], 0.7, y + 0.35, 11.8, 1.55, font_size=9.5, color=TEXT, line_spacing=1.3)

# Sources footnote
add_multi_text(s, [
    "Sources: Market Intelo, Dataintelo, Trustpilot, Reddit, MBW, official pricing pages "
    "(Splice/Loopcloud/Tracklib/Spitfire/NI/Output/Heavyocity/Cinesamples/8Dio/Embertone/"
    "Orchestral Tools/EastWest), Vi-Control, KVR Audio, Gearspace, Attack Magazine, BandLab HC, "
    "MusicRadar, MusicTech, Adobe Podcast, Google Magenta, Samplab/RipX/Waves/Text-to-Sample/Loudly. "
    "Research consolidated August 16–18, 2026."
], 0.55, 6.65, 12.2, 0.35, font_size=7, color=MUTED, line_spacing=1.15)

# ── SAVE ──
prs.save(OUTPUT)
print(f"✓ Slides written to: {OUTPUT}")
print(f"  Size: {os.path.getsize(OUTPUT):,} bytes")
print(f"  Slides: {len(prs.slides)}")
